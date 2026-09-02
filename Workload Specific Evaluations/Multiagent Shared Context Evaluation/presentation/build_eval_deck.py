"""Generate a PPTX deck: persona lenses on agent evaluation, plus the three
cautions that keep persona-partitioned metrics from misleading you.

Run:  python build_eval_deck.py
Output: Evaluating-Agent-Functions.pptx (same directory)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
INK    = RGBColor(0x1A, 0x1F, 0x2B)
NAVY   = RGBColor(0x0F, 0x2A, 0x4A)
BLUE   = RGBColor(0x2D, 0x6C, 0xDF)
TEAL   = RGBColor(0x12, 0x9E, 0x8F)
AMBER  = RGBColor(0xE0, 0x8A, 0x1E)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PURPLE = RGBColor(0x6B, 0x4E, 0xC4)
GREY   = RGBColor(0x5B, 0x64, 0x70)
LIGHT  = RGBColor(0xF2, 0xF5, 0xFA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD6, 0xDE, 0xEA)

SW, SH = Inches(13.333), Inches(7.5)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def _outline(shape, fill, line, weight=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(weight)
    shape.shadow.inherit = False


def _text(host, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6, line_spacing=1.03):
    tf = host.text_frame if hasattr(host, "text_frame") else host
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=WHITE):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _solid(r, color)
    return r


def header(slide, kicker, title, color=NAVY):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.2))
    _solid(band, color)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), SW, Inches(0.06))
    _solid(stripe, BLUE)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.14), Inches(12.1), Inches(1.0))
    _text(tb, [
        [(kicker, 12, RGBColor(0xBF, 0xD0, 0xEA), True)],
        [(title, 27, WHITE, True)],
    ], space_after=2)


def card(slide, x, y, w, h, fill=LIGHT, line=LINE, weight=1.0, rounded=True):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    c = slide.shapes.add_shape(st, x, y, w, h)
    _outline(c, fill, line, weight)
    return c


def chip(slide, x, y, w, text, color, textcolor=WHITE, h=Inches(0.4), size=11.5):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _solid(c, color)
    _text(c, [(text, size, textcolor, True)], align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return c


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

PERSONAS = [
    ("PM / Domain Expert", PURPLE, "Is the output actually good?"),
    ("Application Engineer", BLUE, "Do the contracts hold?"),
    ("Platform Engineer", TEAL, "Is the substrate efficient?"),
    ("Prod Ops / SRE", AMBER, "Is it fast, cheap, stable?"),
]

# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
s = add_slide(prs)
bg(s, NAVY)
b1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.7), SW, Inches(0.12))
_solid(b1, BLUE)
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.6), Inches(3.2))
_text(tb, [
    [("EVALUATING AGENT FUNCTIONS", 15, TEAL, True)],
    [("Different personas, different metrics", 40, WHITE, True)],
    [("Why the same agent looks \u201cgood\u201d or \u201cbroken\u201d depending on who is looking \u2014",
      19, RGBColor(0xC9, 0xD6, 0xEA), False)],
    [("and the three cautions that keep persona dashboards from misleading you", 19,
      RGBColor(0xC9, 0xD6, 0xEA), False)],
], space_after=9)
tb2 = s.shapes.add_textbox(Inches(0.9), Inches(6.05), Inches(11.6), Inches(0.8))
_text(tb2, [("Example: market_trends_analyst  \u2022  Multiagent Shared Context Evaluation",
             13, RGBColor(0x9F, 0xB2, 0xD0), False)])

# ---------------------------------------------------------------------------
# Slide 2 — The persona lens (matrix)
# ---------------------------------------------------------------------------
s = add_slide(prs)
bg(s)
header(s, "THE PERSONA LENS", "The same metric suite, four different questions")
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.7))
_text(tb, [("A metric suite inherits the blind spot of its author. This one was built by a "
            "platform mindset \u2014 great on flow, silent on quality.", 14, GREY, False)])

y = Inches(2.25); h = Inches(2.05)
w = Inches(2.95); gap = Inches(0.13); x0 = Inches(0.55)
mets = [
    "Groundedness\nAnalysis Quality\nState Consistency",
    "Handoff Completeness\nContext Freshness\nContext Utilization\nWrite Accuracy",
    "Redundancy\nContext Compression (CCR)\nC2 Alignment",
    "Read / Write Latency\nCoordination Latency %\nCoordination Token %",
]
for i, ((name, col, q), metblock) in enumerate(zip(PERSONAS, mets)):
    x = x0 + i * (w + gap)
    cd = card(s, x, y, w, h, fill=WHITE, line=col, weight=2.0)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.5))
    _solid(top, col)
    _text(top, [(name, 13, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    tbq = s.shapes.add_textbox(x + Inches(0.18), y + Inches(0.58), w - Inches(0.36), Inches(0.5))
    _text(tbq, [("\u201c" + q + "\u201d", 12.5, col, True)])
    tbm = s.shapes.add_textbox(x + Inches(0.18), y + Inches(1.05), w - Inches(0.36), h - Inches(1.1))
    paras = [[(line, 11.5, INK, False)] for line in metblock.split("\n")]
    _text(tbm, paras, space_after=3)

foot = card(s, Inches(0.55), Inches(4.6), Inches(12.25), Inches(2.15), fill=LIGHT)
_text(foot.text_frame, [
    [("Rows = what \u201cgood job\u201d means (function questions):", 14, NAVY, True)],
    [("input fidelity  \u2022  core task quality  \u2022  output fidelity  \u2022  resource efficiency", 13, INK, False)],
    [("Columns = who leads on each row, and when. Same underlying scores \u2014 different projection.", 13, INK, False)],
    [("The gap we found: the \u201ccore task quality\u201d row was empty because no PM/domain persona "
      "was in the room when the suite was authored.", 13, PURPLE, True)],
], space_after=6)
foot.text_frame.margin_left = Inches(0.3); foot.text_frame.margin_top = Inches(0.18)

# ---------------------------------------------------------------------------
# Slide 3 — Cautions intro
# ---------------------------------------------------------------------------
s = add_slide(prs)
bg(s)
header(s, "THE CATCH", "Partitioning by persona helps triage \u2014 but three things bite")
items = [
    ("1", RED, "Failures hide in the seams", "Green for one persona, red for another \u2014 on the same event."),
    ("2", AMBER, "Same metric, different meaning", "One number, several interpretations and thresholds."),
    ("3", PURPLE, "Persona \u2260 lifecycle stage", "Who is looking is a different axis from when."),
]
y = Inches(1.9); h = Inches(1.4); w = Inches(11.9); x = Inches(0.7)
for i, (n, col, t, b) in enumerate(items):
    yy = y + Emu(int(h) * i) + Inches(0.15) * i
    cd = card(s, x, yy, w, h, fill=WHITE, line=col, weight=2.0)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), yy + Inches(0.35), Inches(0.7), Inches(0.7))
    _solid(circ, col)
    _text(circ, [(n, 24, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    tbc = s.shapes.add_textbox(x + Inches(1.3), yy + Inches(0.22), w - Inches(1.6), h - Inches(0.4))
    _text(tbc, [
        [(t, 21, INK, True)],
        [(b, 14, GREY, False)],
    ], space_after=5, anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------------------
# Helper for a full caution slide
# ---------------------------------------------------------------------------

def caution_slide(kicker, title, color, left_pairs, callout_label, callout_text,
                  callout_color):
    s = add_slide(prs)
    bg(s)
    header(s, kicker, title, color=NAVY)
    tag = chip(s, Inches(0.6), Inches(1.45), Inches(2.2), kicker, color, h=Inches(0.42), size=12)
    # left: the mechanism (two contrasting cards)
    y = Inches(2.15); h = Inches(1.55); w = Inches(6.0); x = Inches(0.6)
    for i, (head, body, hc) in enumerate(left_pairs):
        yy = y + Emu(int(h) * i) + Inches(0.2) * i
        cd = card(s, x, yy, w, h, fill=WHITE, line=hc, weight=1.75)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, yy, Inches(0.13), h)
        _solid(bar, hc)
        tbc = s.shapes.add_textbox(x + Inches(0.32), yy + Inches(0.16), w - Inches(0.55), h - Inches(0.3))
        _text(tbc, [
            [(head, 15, hc, True)],
            [(body, 13, INK, False)],
        ], space_after=5)
    # right: callout
    cx = Inches(6.9); cw = Inches(5.85)
    cc = card(s, cx, Inches(2.15), cw, Inches(4.3), fill=LIGHT, line=callout_color, weight=2.0)
    _text(cc.text_frame, callout_text, space_after=9)
    cc.text_frame.margin_left = Inches(0.32)
    cc.text_frame.margin_right = Inches(0.28)
    cc.text_frame.margin_top = Inches(0.28)
    lbl = chip(s, cx + Inches(0.32), Inches(2.35), Inches(3.4), callout_label, callout_color,
               h=Inches(0.42), size=12)
    return s

# ---------------------------------------------------------------------------
# Slide (new) — The evaluation dataset: how you feed the pipeline
# ---------------------------------------------------------------------------
s = add_slide(prs)
bg(s)
header(s, "THE DATASET", "What you evaluate against: grounded research briefs")
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.62))
_text(tb, [("You cannot judge \u201cgood analysis\u201d without a reference. Step 1 of the loop is a "
            "dataset: briefs + the verifiable facts each one states.", 14, GREY, False)])

# Pipeline strip
py = Inches(2.2); pw = Inches(2.55); ph = Inches(0.85); pgap = Inches(0.55); px0 = Inches(1.05)
pipeline = [
    ("Research brief", PURPLE),
    ("market_trends", BLUE),
    ("customer_insights", TEAL),
    ("strategy_synth", NAVY),
]
for i, (name, col) in enumerate(pipeline):
    x = px0 + i * (pw + pgap)
    c = card(s, x, py, pw, ph, fill=WHITE, line=col, weight=2.0)
    _text(c, [(name, 13, col, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    if i < len(pipeline) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + pw + Emu(1000), py + Inches(0.28),
                                pgap - Emu(2000), Inches(0.3))
        _solid(ar, AMBER)

# Left: the JSONL row shape
ly = Inches(3.45); lw = Inches(6.0); lh = Inches(2.75); lx = Inches(0.6)
lc = card(s, lx, ly, lw, lh, fill=RGBColor(0x1B, 0x22, 0x2E))
_text(lc.text_frame, [
    [("datasets/research_briefs.jsonl", 13, RGBColor(0x7F, 0xD0, 0xC8), True)],
    [("{", 12, WHITE, False)],
    [("  \"id\": \"brief-01-fastcasual\",", 12, RGBColor(0xD6, 0xDE, 0xEA), False)],
    [("  \"turns\": [ \"Analyse the fast-casual market\u2026\" ],", 12, RGBColor(0xD6, 0xDE, 0xEA), False)],
    [("  \"facts\": { market_size: $60B, growth: 8%,", 12, RGBColor(0x9C, 0xD1, 0xFF), False)],
    [("            players: [Chipotle, Panera, Sweetgreen] },", 12, RGBColor(0x9C, 0xD1, 0xFF), False)],
    [("  \"premise_valid\": true", 12, RGBColor(0xD6, 0xDE, 0xEA), False)],
    [("}", 12, WHITE, False)],
], space_after=4)
lc.text_frame.margin_left = Inches(0.28); lc.text_frame.margin_top = Inches(0.2)

# Right: how facts are used
ry = Inches(3.45); rw = Inches(5.85); rx = Inches(6.9)
rc = card(s, rx, ry, rw, lh, fill=LIGHT, line=PURPLE, weight=1.5)
_text(rc.text_frame, [
    [("facts  \u2192  source_facts  \u2192  Groundedness judge", 14, PURPLE, True)],
    [("The stated facts become the ground-truth reference. The judge checks whether "
      "each agent\u2019s output stays faithful to them \u2014 not just internally consistent.", 13, INK, False)],
    [("", 5, GREY, False)],
    [("Option A (now):", 12.5, NAVY, True), (" briefs + verifiable facts, no gold answers.", 12.5, INK, False)],
    [("Option B (later):", 12.5, GREY, True), (" add per-stage \u201cgood answer\u201d references.", 12.5, GREY, False)],
], space_after=6)
rc.text_frame.margin_left = Inches(0.3); rc.text_frame.margin_top = Inches(0.2)

# Trap chip: flawed-premise row
trap = card(s, Inches(6.9), Inches(6.35), Inches(5.85), Inches(0.75),
            fill=RGBColor(0xFD, 0xF3, 0xE2), line=AMBER, weight=1.5)
_text(trap.text_frame, [
    [("\u26a0 One brief has a flawed premise", 12.5, AMBER, True),
     (" \u2014 a good analyst should flag it, not restate the fake numbers.", 12, INK, False)],
], space_after=0, anchor=MSO_ANCHOR.MIDDLE)
trap.text_frame.margin_left = Inches(0.24)

# ---------------------------------------------------------------------------
# Slide (new) — Validating the judge: human-in-the-loop
# ---------------------------------------------------------------------------
s = add_slide(prs)
bg(s)
header(s, "TRUST THE JUDGE", "Who validates the validator? A human-in-the-loop benchmark")
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.42), Inches(12.1), Inches(0.6))
_text(tb, [("An LLM grades the pipeline. But nothing grades the LLM \u2014 until a human "
            "labels ground truth and we measure how often the judge agrees.", 14, GREY, False)])

# 5-step workflow strip
steps = [
    ("1  RUN", BLUE, "Run every brief through the pipeline (live)."),
    ("2  CAPTURE", TEAL, "Dump each output to a review file with its expected behavior."),
    ("3  LABEL", PURPLE, "Human reads output + rule, marks PASS / FAIL + one-line reason."),
    ("4  SCORE", AMBER, "Run the judge on the same cases; compare to the human labels."),
    ("5  FOLD BACK", NAVY, "Labels become the benchmark; re-run as the judge/prompts change."),
]
w = Inches(2.28); h = Inches(1.9); gap = Inches(0.16); x0 = Inches(0.55); y = Inches(2.35)
for i, (label, col, body) in enumerate(steps):
    x = x0 + i * (w + gap)
    card(s, x, y, w, h, fill=WHITE, line=col, weight=2.0)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.46))
    _solid(top, col)
    _text(top, [(label, 12.5, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    tbc = s.shapes.add_textbox(x + Inches(0.18), y + Inches(0.56), w - Inches(0.36), h - Inches(0.66))
    _text(tbc, [(body, 12, INK, False)], space_after=0)
    if i < len(steps) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w + Emu(500), y + Inches(0.7),
                                gap - Emu(1000), Inches(0.4))
        _solid(ar, AMBER)

# The metric that matters + ground-truth note (two side-by-side cards)
lc = card(s, Inches(0.55), Inches(4.65), Inches(6.0), Inches(2.05), fill=RGBColor(0xFD, 0xF3, 0xE2),
          line=RED, weight=2.0)
_text(lc.text_frame, [
    [("Watch the false positives", 15, RED, True)],
    [("The dangerous error is the judge passing an output a human failed \u2014 a "
      "silent quality miss. Track those, not just overall accuracy.", 13.5, INK, False)],
    [("", 5, GREY, False)],
    [("Confusion matrix, human vs judge:", 12.5, NAVY, True),
     ("  pass/pass, fail/fail = agree;  fail\u2192pass = false positive.", 12.5, INK, False)],
], space_after=6)
lc.text_frame.margin_left = Inches(0.3); lc.text_frame.margin_top = Inches(0.2)

rc = card(s, Inches(6.75), Inches(4.65), Inches(6.0), Inches(2.05), fill=LIGHT, line=PURPLE, weight=1.5)
_text(rc.text_frame, [
    [("The human is the ground truth", 15, PURPLE, True)],
    [("You grade against the case\u2019s expected behavior \u2014 not \u201cis this good "
      "analysis?\u201d Cite specific evidence; when unsure, default to FAIL.", 13.5, INK, False)],
    [("", 5, GREY, False)],
    [("Seed labels are synthetic", 12.5, NAVY, True),
     (" until replaced with human-labeled real outputs.", 12.5, INK, False)],
], space_after=6)
rc.text_frame.margin_left = Inches(0.3); rc.text_frame.margin_top = Inches(0.2)

# ---------------------------------------------------------------------------
# Slide 4 — Caution 1: seams
# ---------------------------------------------------------------------------
caution_slide(
    "CAUTION 1", "Failures hide in the seams between personas", RED,
    left_pairs=[
        ("Engineer's view: Write Accuracy = 5/5",
         "\u201cThe output is consistent with the input it was given.\u201d  \u2192 looks healthy.", TEAL),
        ("PM's view: Groundedness = 2/5",
         "\u201cThe output does not match reality \u2014 the input premise was wrong.\u201d  \u2192 broken.", RED),
    ],
    callout_label="THE TRAP",
    callout_text=[
        [("Same agent. Same turn. Two personas, opposite verdicts.", 16, INK, True)],
        [("If each persona reads only their own column, nobody owns the seam \u2014 "
          "and the seam is exactly where the failure lives.", 14, GREY, False)],
        [("", 6, GREY, False)],
        [("Personas are lenses, not owners.", 15, RED, True)],
        [("Everyone gets a home dashboard, but quality gates require a "
          "cross-persona read. Partitioning improves triage; taken as strict "
          "ownership it worsens detection.", 13.5, INK, False)],
    ],
    callout_color=RED,
)

# ---------------------------------------------------------------------------
# Slide 5 — Caution 2: same metric, different meaning
# ---------------------------------------------------------------------------
caution_slide(
    "CAUTION 2", "The same metric means different things", AMBER,
    left_pairs=[
        ("Coordination Token %",
         "SRE reads it as COST. Engineer reads it as context bloat / prompt quality.", AMBER),
        ("Latency",
         "Ops reads it as an SLA. PM reads it as user-perceived responsiveness.", BLUE),
    ],
    callout_label="THE TRAP",
    callout_text=[
        [("One number \u2260 one meaning.", 16, INK, True)],
        [("The metric is shared; the interpretation and the threshold are not.", 14, GREY, False)],
        [("", 6, GREY, False)],
        [("A persona view = metric + interpretation + threshold", 15, AMBER, True)],
        [("\u2026not just a filtered column list. Don\u2019t assume a shared "
          "definition just because you share a column header.", 13.5, INK, False)],
    ],
    callout_color=AMBER,
)

# ---------------------------------------------------------------------------
# Slide 6 — Caution 3: persona vs stage
# ---------------------------------------------------------------------------
caution_slide(
    "CAUTION 3", "Persona is a different axis from lifecycle stage", PURPLE,
    left_pairs=[
        ("WHO  (persona)",
         "PM / Domain \u2022 App Engineer \u2022 Platform \u2022 SRE", PURPLE),
        ("WHEN  (stage)",
         "Design \u2022 Code \u2022 Deploy \u2022 Operate", BLUE),
    ],
    callout_label="THE TRAP",
    callout_text=[
        [("They correlate \u2014 but they are not the same axis.", 16, INK, True)],
        [("Groundedness is not \u201cdesign-only.\u201d It is a design-time gate "
          "AND a production canary (e.g. drift after a model swap).", 14, GREY, False)],
        [("", 6, GREY, False)],
        [("Keep two axes: who + when.", 15, PURPLE, True)],
        [("The same metric recurs at multiple stages with different urgency. "
          "Collapsing the axes makes coverage look complete when it isn\u2019t.", 13.5, INK, False)],
    ],
    callout_color=PURPLE,
)

# ---------------------------------------------------------------------------
# Slide 7 — Takeaway
# ---------------------------------------------------------------------------
s = add_slide(prs)
bg(s, NAVY)
tb = s.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0))
_text(tb, [("So what do you do with this?", 30, WHITE, True)])
stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.55), Inches(3.0), Inches(0.06))
_solid(stripe, TEAL)
points = [
    ("Give every persona a home dashboard.", "Faster triage \u2014 SRE sees latency, PM sees quality, no wading."),
    ("Gate releases on a cross-persona read.", "The dangerous failures live in the seams; require more than one lens to pass."),
    ("Tag each metric with {persona, stage}.", "Same judge scores, rendered as different projections. Two axes, not one."),
    ("Treat \u201cgood job\u201d as a definition you sharpen.", "Measure \u2192 find the empty row \u2192 add the missing persona\u2019s metric \u2192 re-run."),
]
y = Inches(1.95)
for i, (h, b) in enumerate(points):
    yy = y + Inches(1.2) * i
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), yy + Inches(0.04), Inches(0.34), Inches(0.34))
    _solid(dot, TEAL)
    _text(dot, [(str(i + 1), 14, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    tbp = s.shapes.add_textbox(Inches(1.5), yy, Inches(11.0), Inches(1.05))
    _text(tbp, [
        [(h, 18, WHITE, True)],
        [(b, 14, RGBColor(0xC9, 0xD6, 0xEA), False)],
    ], space_after=3)

out = "Evaluating-Agent-Functions.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
